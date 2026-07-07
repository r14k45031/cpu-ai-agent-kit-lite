# -*- coding: utf-8 -*-
"""共用模組：呼叫本機 Ollama、讀取文件、OCR、輸出 Word。全程離線。"""
import ctypes
import json
import os
import shutil
import subprocess
import sys
import tempfile
import urllib.request
from pathlib import Path

OLLAMA = "http://127.0.0.1:11434"
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff"}


def _post_json(url, payload):
    req = urllib.request.Request(
        url,
        data=json.dumps(payload).encode("utf-8"),
        headers={"Content-Type": "application/json"},
    )
    return urllib.request.urlopen(req)


def list_models():
    try:
        with urllib.request.urlopen(OLLAMA + "/api/tags", timeout=5) as r:
            return [m["name"] for m in json.load(r).get("models", [])]
    except OSError:
        raise SystemExit("連不上 Ollama 服務。請先開啟 Ollama（或重新執行 install.bat）。")


def total_ram_gb():
    class MEMORYSTATUSEX(ctypes.Structure):
        _fields_ = [("dwLength", ctypes.c_ulong), ("dwMemoryLoad", ctypes.c_ulong)] + [
            (n, ctypes.c_ulonglong)
            for n in (
                "ullTotalPhys", "ullAvailPhys", "ullTotalPageFile",
                "ullAvailPageFile", "ullTotalVirtual", "ullAvailVirtual",
                "ullAvailExtendedVirtual",
            )
        ]

    m = MEMORYSTATUSEX()
    m.dwLength = ctypes.sizeof(MEMORYSTATUSEX)
    ctypes.windll.kernel32.GlobalMemoryStatusEx(ctypes.byref(m))
    return m.ullTotalPhys / 2**30


_model_cache = None


def pick_model():
    """記憶體夠就用 qwen3-30b，否則退回 qwen3-4b。"""
    global _model_cache
    if _model_cache:
        return _model_cache
    models = [m for m in list_models() if m.startswith("qwen3")]
    if not models:
        raise SystemExit("找不到 qwen3 模型，請先執行 install.bat 完成安裝。")
    big = [m for m in models if "30b" in m]
    small = [m for m in models if "4b" in m]
    if big and total_ram_gb() >= 28:
        _model_cache = big[0]
    else:
        _model_cache = (small or models)[0]
    return _model_cache


def pick_vision_model():
    """找已安裝的視覺模型（qwen3-vl）。"""
    vl = [m for m in list_models() if m.startswith("qwen3-vl")]
    if not vl:
        raise SystemExit("找不到視覺模型 qwen3-vl，請先執行 install.bat 完成安裝。")
    return vl[0]


def chat(prompt, model=None, image_path=None):
    """對本機模型送出提示，串流接收，回傳完整回答。可附一張圖片。"""
    msg = {"role": "user", "content": prompt}
    if image_path:
        import base64

        msg["images"] = [base64.b64encode(Path(image_path).read_bytes()).decode()]
        model = model or pick_vision_model()
    else:
        model = model or pick_model()
    payload = {
        "model": model,
        "messages": [msg],
        "stream": True,
    }
    pieces = []
    n = 0
    with _post_json(OLLAMA + "/api/chat", payload) as r:
        for line in r:
            d = json.loads(line)
            piece = d.get("message", {}).get("content", "")
            if piece:
                pieces.append(piece)
                n += 1
                if n % 30 == 0:
                    print(".", end="", flush=True)
    if n >= 30:
        print()
    return "".join(pieces).strip()


def chunk_text(text, size=6000):
    """依段落切塊，每塊約 size 字。"""
    paras = text.split("\n")
    chunks, buf, length = [], [], 0
    for p in paras:
        if length + len(p) > size and buf:
            chunks.append("\n".join(buf))
            buf, length = [], 0
        buf.append(p)
        length += len(p) + 1
    if buf:
        chunks.append("\n".join(buf))
    return chunks or [""]


def find_tesseract():
    for c in (
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        os.path.expandvars(r"%LOCALAPPDATA%\Programs\Tesseract-OCR\tesseract.exe"),
        shutil.which("tesseract"),
    ):
        if c and Path(c).is_file():
            return c
    return None


def ocr_pdf_text(path):
    """掃描型 PDF → 逐頁轉圖 → Tesseract 辨識（繁中＋英文）。"""
    import pymupdf

    tess = find_tesseract()
    if not tess:
        raise SystemExit("未安裝 Tesseract OCR，無法辨識掃描檔。請重新執行 install.bat。")
    doc = pymupdf.open(str(path))
    parts = []
    with tempfile.TemporaryDirectory() as td:
        for i, page in enumerate(doc, 1):
            img = os.path.join(td, f"p{i}.png")
            page.get_pixmap(dpi=300).save(img)
            out = os.path.join(td, f"p{i}")
            subprocess.run(
                [tess, img, out, "-l", "chi_tra+eng"],
                check=True, capture_output=True,
            )
            parts.append(Path(out + ".txt").read_text(encoding="utf-8", errors="ignore"))
            print(f"  OCR 第 {i}/{len(doc)} 頁完成")
    return "\n".join(parts)


def ocr_image_text(path):
    """單張圖片（PNG/JPG/TIF/BMP）→ Tesseract 辨識（繁中＋英文）。"""
    tess = find_tesseract()
    if not tess:
        raise SystemExit("未安裝 Tesseract OCR，無法辨識圖片。請重新執行 install.bat。")
    with tempfile.TemporaryDirectory() as td:
        out = os.path.join(td, "out")
        subprocess.run(
            [tess, str(path), out, "-l", "chi_tra+eng"],
            check=True, capture_output=True,
        )
        return Path(out + ".txt").read_text(encoding="utf-8", errors="ignore")


def read_pptx_text(path):
    """讀出 PowerPoint 每張投影片的文字。"""
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [sh.text for sh in slide.shapes
                 if sh.has_text_frame and sh.text.strip()]
        parts.append(f"【第 {i} 張投影片】\n" + "\n".join(texts))
    return "\n\n".join(parts)


def read_excel_text(path, max_rows=200):
    """把試算表轉成文字表格（每個工作表最多 max_rows 列）。"""
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        lines = [f"【工作表：{ws.title}】（{ws.max_row} 列 × {ws.max_column} 欄）"]
        for r, row in enumerate(ws.iter_rows(values_only=True), 1):
            if r > max_rows:
                lines.append(f"…（其餘 {ws.max_row - max_rows} 列略）")
                break
            lines.append("\t".join("" if c is None else str(c) for c in row))
        parts.append("\n".join(lines))
    wb.close()
    return "\n\n".join(parts)


def read_document(path):
    """讀出 PDF / Word / PowerPoint / Excel / 圖片 / 純文字的文字內容；
    掃描檔與圖片自動走 OCR。"""
    p = Path(path)
    ext = p.suffix.lower()
    if ext in IMAGE_EXTS:
        print("  圖片檔，使用 OCR 辨識文字...")
        return ocr_image_text(p)
    if ext == ".pptx":
        return read_pptx_text(p)
    if ext in (".xlsx", ".xlsm"):
        return read_excel_text(p)
    if ext == ".pdf":
        import pymupdf

        doc = pymupdf.open(str(p))
        text = "\n".join(page.get_text() for page in doc)
        if len(text.strip()) < 80 and find_tesseract():
            print("  這份 PDF 沒有文字層，改用 OCR 辨識（較花時間）...")
            text = ocr_pdf_text(p)
        return text
    if ext == ".docx":
        import docx

        d = docx.Document(str(p))
        return "\n".join(par.text for par in d.paragraphs)
    return p.read_text(encoding="utf-8", errors="ignore")


def write_docx(path, title, body):
    """把純文字（可含 # 標題、- 條列的簡易 Markdown）寫成 Word 檔。"""
    import docx

    d = docx.Document()
    if title:
        d.add_heading(title, level=1)
    for line in body.splitlines():
        s = line.strip()
        if not s:
            continue
        if s.startswith("### "):
            d.add_heading(s[4:], level=3)
        elif s.startswith("## "):
            d.add_heading(s[3:], level=2)
        elif s.startswith("# "):
            d.add_heading(s[2:], level=1)
        elif s.startswith(("- ", "* ")):
            d.add_paragraph(s[2:], style="List Bullet")
        else:
            d.add_paragraph(s)
    d.save(str(path))


if __name__ == "__main__":
    print("這是共用模組，請執行 summarize.py / ocr_pdf.py / polish_docx.py，"
          "或把檔案拖到 tools 資料夾中的 .bat 上。")
