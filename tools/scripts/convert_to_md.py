# -*- coding: utf-8 -*-
"""文件轉 Markdown：把 Word/PDF/PPT/Excel/圖片統一轉成 .md 純文字格式。

確定性轉換、不經過 AI——速度快、內容不失真，Ollama 沒開也能用。
適合：整理資料餵知識庫、純文字存檔、進版本控制。

用法：convert_to_md.py <檔案或資料夾路徑>
輸出：與原檔同資料夾的同名 .md
"""
import os
import shutil
import subprocess
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
import ai_common as ai

EXTS = {".docx", ".pdf", ".pptx", ".xlsx", ".xlsm", ".csv", ".txt"} | ai.IMAGE_EXTS


def find_pandoc():
    for c in (os.path.expandvars(r"%LOCALAPPDATA%\Pandoc\pandoc.exe"),
              shutil.which("pandoc")):
        if c and Path(c).is_file():
            return c
    return None


def cell(v):
    return str(v).replace("|", "\\|").replace("\n", " ").strip()


def docx_to_md(f):
    pandoc = find_pandoc()
    if pandoc:
        r = subprocess.run(
            [pandoc, "-f", "docx", "-t", "gfm", "--wrap=none", str(f)],
            capture_output=True,
        )
        if r.returncode == 0 and r.stdout.strip():
            return r.stdout.decode("utf-8", errors="ignore")
    # 後備：沒有 pandoc 時用 python-docx 手動轉
    import docx

    d = docx.Document(str(f))
    lines = []
    for p in d.paragraphs:
        t = p.text.strip()
        if not t:
            continue
        style = p.style.name or ""
        if style.startswith("Heading"):
            try:
                lvl = int(style.split()[-1])
            except ValueError:
                lvl = 2
            lines.append("#" * min(lvl, 6) + " " + t)
        elif style.startswith("List"):
            lines.append("- " + t)
        else:
            lines.append(t)
        lines.append("")
    for tbl in d.tables:
        rows = [[cell(c.text) for c in r.cells] for r in tbl.rows]
        if not rows:
            continue
        lines.append("| " + " | ".join(rows[0]) + " |")
        lines.append("|" + "---|" * len(rows[0]))
        for r in rows[1:]:
            lines.append("| " + " | ".join(r) + " |")
        lines.append("")
    return "\n".join(lines)


def pdf_to_md(f):
    try:
        import pymupdf4llm

        md = pymupdf4llm.to_markdown(str(f), show_progress=False)
        if len(md.strip()) >= 80:
            return md
    except Exception as e:
        print(f"  pymupdf4llm 轉換失敗（{e}），改用純文字擷取")
    return ai.read_document(f)  # 純文字；沒有文字層會自動 OCR


def pptx_to_md(f):
    from pptx import Presentation

    prs = Presentation(str(f))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape is not None else ""
        title_id = title_shape.shape_id if title_shape is not None else None
        lines.append(f"## 投影片 {i}" + (f"：{title}" if title else ""))
        for sh in slide.shapes:
            if not sh.has_text_frame or sh.shape_id == title_id:
                continue
            for para in sh.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    lines.append("- " + t)
        lines.append("")
    return "\n".join(lines)


def sheet_to_md(rows_iter, total_rows, max_rows=500):
    lines = []
    first = True
    for n, row in enumerate(rows_iter, 1):
        if n > max_rows:
            lines.append(f"（其餘 {total_rows - max_rows} 列略）")
            break
        cells = ["" if c is None else cell(c) for c in row]
        lines.append("| " + " | ".join(cells) + " |")
        if first:
            lines.append("|" + "---|" * len(cells))
            first = False
    return lines


def xlsx_to_md(f):
    import openpyxl

    wb = openpyxl.load_workbook(str(f), read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append(f"## 工作表：{ws.title}")
        lines += sheet_to_md(ws.iter_rows(values_only=True), ws.max_row)
        lines.append("")
    wb.close()
    return "\n".join(lines)


def csv_to_md(f):
    import csv

    with open(f, newline="", encoding="utf-8-sig", errors="ignore") as fh:
        rows = list(csv.reader(fh))
    return "\n".join(sheet_to_md(iter(rows), len(rows)))


def convert(f):
    ext = f.suffix.lower()
    if ext == ".docx":
        md = docx_to_md(f)
    elif ext == ".pdf":
        md = pdf_to_md(f)
    elif ext == ".pptx":
        md = pptx_to_md(f)
    elif ext in (".xlsx", ".xlsm"):
        md = xlsx_to_md(f)
    elif ext == ".csv":
        md = csv_to_md(f)
    elif ext in ai.IMAGE_EXTS:
        print("  圖片檔，OCR 辨識中...")
        md = ai.ocr_image_text(f)
    else:
        md = f.read_text(encoding="utf-8", errors="ignore")
    out = f.with_suffix(".md")
    if out.exists():
        out = f.with_name(f.stem + "_轉換.md")
    out.write_text(md.strip() + "\n", encoding="utf-8")
    return out


def main():
    if len(sys.argv) < 2:
        raise SystemExit("用法：把檔案或資料夾拖到「拖入文件轉MD.bat」上")
    target = Path(sys.argv[1])
    if target.is_file():
        files = [target]
    else:
        files = sorted(p for p in target.iterdir()
                       if p.is_file() and p.suffix.lower() in EXTS)
    if not files:
        raise SystemExit("找不到可轉換的檔案（支援 docx/pdf/pptx/xlsx/csv/txt/圖片）")
    print(f"共 {len(files)} 個檔案")
    for f in files:
        print(f"轉換：{f.name}")
        try:
            out = convert(f)
            print(f"  → {out.name}")
        except Exception as e:
            print(f"  失敗：{e}")
    print("全部完成。")


if __name__ == "__main__":
    main()
