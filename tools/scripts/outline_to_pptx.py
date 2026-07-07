# -*- coding: utf-8 -*-
"""文件變簡報：讀入文件（Word/PDF/txt/md），AI 整理大綱後自動產生 PowerPoint。

用法：outline_to_pptx.py <文件路徑>
輸出：<原檔名>_簡報.pptx
"""
import json
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
import ai_common as ai

PROMPT = """請把以下內容整理成簡報大綱，以 JSON 陣列輸出，格式：
[{"title": "投影片標題", "bullets": ["要點一", "要點二"]}]
規則：第一張是整份簡報的標題頁（bullets 放 2-3 句副標或重點）；
共 6～12 張；每張最多 5 個要點，每個要點不超過 25 字；
只輸出 JSON，不要其他文字。

內容：
"""


def parse_json_array(text):
    m = re.search(r"\[.*\]", text, re.DOTALL)
    if not m:
        raise ValueError("回應中找不到 JSON")
    return json.loads(m.group(0))


def build_pptx(slides, out_path):
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    body_layout = prs.slide_layouts[1]
    for i, s in enumerate(slides):
        layout = title_layout if i == 0 else body_layout
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = str(s.get("title", ""))[:60]
        body = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:
                body = ph
                break
        if body is None:
            continue
        tf = body.text_frame
        bullets = [str(b) for b in s.get("bullets", []) if str(b).strip()]
        for j, b in enumerate(bullets[:6]):
            para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            para.text = b
            para.font.size = Pt(20 if i else 18)
    prs.save(str(out_path))


def main():
    if len(sys.argv) < 2:
        raise SystemExit("用法：把文件拖到「拖入文件變簡報.bat」上")
    f = Path(sys.argv[1])
    if f.suffix.lower() not in (".docx", ".pdf", ".txt", ".md"):
        raise SystemExit("請拖入 .docx / .pdf / .txt / .md 檔案")
    model = ai.pick_model()
    print(f"讀取：{f.name}（模型：{model}）")
    text = ai.read_document(f)
    if len(text.strip()) < 40:
        raise SystemExit("擷取不到足夠文字。")
    text = "\n".join(ai.chunk_text(text, size=6000)[:1])
    print("AI 整理簡報大綱中...")
    slides = None
    for attempt in (1, 2):
        try:
            slides = parse_json_array(ai.chat(PROMPT + text, model=model))
            break
        except (ValueError, json.JSONDecodeError):
            if attempt == 1:
                print("  大綱格式不對，重試一次...")
    if not slides:
        raise SystemExit("AI 兩次都沒有產出有效大綱，請再試一次或換一份文件。")
    out = f.with_name(f.stem + "_簡報.pptx")
    build_pptx(slides, out)
    print(f"完成：{len(slides)} 張投影片。已輸出：{out.name}")
    print("（開啟後可在 PowerPoint 裡套用任何佈景主題美化）")


if __name__ == "__main__":
    main()
